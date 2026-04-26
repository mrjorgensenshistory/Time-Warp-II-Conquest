"""Make PPTX-rendered buttons clickable on all 7 character HTMLs.

Same pattern as analyzer.html fix:
1. Use python-pptx to extract every button text-box's position (left/top/w/h
   as % of slide dimensions) from each source PPTX.
2. Match buttons by TEXT (case-insensitive prefix) to entries in slide.buttons[]
   in the existing SLIDES JSON.
3. Add `pptx_pos: {left, top, w, h}` field to each matched button.
4. Inject CSS + JS engine extension that:
   - Renders matched buttons as positioned-absolute hotspots over the bg_image
     at their PPTX coordinates (with always-visible faint gold dashed border +
     pulsing glow so kids see them clearly).
   - HIDES the default bottom-row engine buttons on slides where ALL buttons
     have positions (so no redundancy).
   - Falls back to bottom-row rendering if positions weren't matched (so we
     never lose interactivity).
"""

import os
import re
import json
from pptx import Presentation

ROOT = os.path.dirname(os.path.abspath(__file__))

PPTX_MAP = {
    "columbus.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
    "cortereal.html": r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Corte-Real) FINAL.pptx",
    "dagama.html":    r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Da Gama) FINAL.pptx",
    "drake.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Drake) FINAL.pptx",
    "magellan.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Magellan) FINAL.pptx",
    "narvaez.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Narvaez) FINAL.pptx",
    "raleigh.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Raleigh) FINAL.pptx",
    # Index uses the Tutorial PPTX
    "index.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
}

SLIDES_RE = re.compile(r"const SLIDES = (\[.*?\]);", re.DOTALL)
ENGINE_MARKER = "<!-- TW2-PPTX-BUTTON-OVERLAY-V1 -->"

ENGINE_BLOCK = """
<!-- TW2-PPTX-BUTTON-OVERLAY-V1 -->
<style>
.pptx-hotspot {
  position: absolute;
  z-index: 8;
  cursor: pointer;
  border: 2px dashed rgba(201, 161, 74, 0.55);
  background: rgba(201, 161, 74, 0.06);
  transition: all 0.15s ease;
  border-radius: 4px;
  box-shadow: 0 0 12px rgba(201, 161, 74, 0.15) inset;
}
.pptx-hotspot:hover {
  border-color: rgba(255, 215, 0, 1);
  border-style: solid;
  background: rgba(201, 161, 74, 0.2);
  box-shadow: 0 0 18px rgba(255, 215, 0, 0.4);
  transform: scale(1.02);
}
.pptx-hotspot:active { background: rgba(201, 161, 74, 0.4); }
@keyframes pptx-hotspot-pulse {
  0%, 100% { box-shadow: 0 0 12px rgba(201,161,74,0.12) inset, 0 0 6px rgba(201,161,74,0.3); }
  50%      { box-shadow: 0 0 12px rgba(201,161,74,0.25) inset, 0 0 14px rgba(255,215,0,0.55); }
}
.pptx-hotspot { animation: pptx-hotspot-pulse 2.4s ease-in-out infinite; }

/* When a slide has ALL buttons positioned, hide the redundant bottom row */
.slide.has-positioned-buttons .buttons {
  display: none !important;
}
</style>
<script>
(function() {
  if (typeof window.goToSlide !== 'function') return;
  // Wrap goToSlide to inject pptx-positioned hotspots after default render
  var prev = window.goToSlide;
  window.goToSlide = function(id) {
    prev(id);
    if (id === 'HOME' || (id || '').endsWith('.html')) return;
    if (typeof SLIDES === 'undefined') return;
    var slide = SLIDES.find(function(s){ return s.id === id; });
    if (!slide || !slide.buttons || !slide.buttons.length) return;

    setTimeout(function() {
      var active = document.querySelector('.slide.active');
      if (!active) return;
      var positionedCount = 0;
      slide.buttons.forEach(function(btn) {
        if (!btn.pptx_pos) return;
        var hs = document.createElement('div');
        hs.className = 'pptx-hotspot';
        hs.style.left = btn.pptx_pos.left + '%';
        hs.style.top = btn.pptx_pos.top + '%';
        hs.style.width = btn.pptx_pos.w + '%';
        hs.style.height = btn.pptx_pos.h + '%';
        hs.title = btn.text;
        hs.onclick = function() { window.goToSlide(btn.target); };
        active.appendChild(hs);
        positionedCount++;
      });
      // Hide bottom row only if every button is positioned
      if (positionedCount === slide.buttons.length && positionedCount > 0) {
        active.classList.add('has-positioned-buttons');
      }
    }, 60);
  };
})();
</script>
"""


def extract_pptx_button_positions(pptx_path):
    """Returns list of {page_idx (1-based), text_lower, left%, top%, w%, h%}"""
    p = Presentation(pptx_path)
    sw, sh = p.slide_width, p.slide_height
    out = []
    for i, slide in enumerate(p.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text = shape.text_frame.text.strip()
            if not text: continue
            # Skip very large shapes (likely body/title containers, not buttons)
            try:
                w_pct = shape.width / sw * 100
                h_pct = shape.height / sh * 100
                # Heuristic: a "button" is < 70% of slide width and < 25% of slide height
                if w_pct > 80 or h_pct > 35:
                    continue
                out.append({
                    "page": i,
                    "text": text,
                    "text_lower": text.lower(),
                    "left": round(shape.left / sw * 100, 1),
                    "top": round(shape.top / sh * 100, 1),
                    "w": round(w_pct, 1),
                    "h": round(h_pct, 1),
                })
            except Exception:
                pass
    return out


def match_button(slide_id, slides, pptx_buttons):
    """For each button in the slide, find a position match by text prefix."""
    # slide_id format: slide_N -> page (N+1)
    try:
        page_idx = int(slide_id.split("_")[1]) + 1
    except Exception:
        return 0

    # Get pptx button shapes for this page
    page_btns = [b for b in pptx_buttons if b["page"] == page_idx]
    if not page_btns:
        return 0

    slide = next((s for s in slides if s.get("id") == slide_id), None)
    if not slide or not slide.get("buttons"):
        return 0

    n_matched = 0
    for btn in slide["buttons"]:
        btn_text = (btn.get("text") or "").lower().strip()
        if not btn_text:
            continue
        # Match: button text appears in PPTX text (button label might be truncated in JSON)
        # Use first 15 chars as a reasonable prefix
        prefix = btn_text[:15]
        match = next((p for p in page_btns
                      if prefix in p["text_lower"] or p["text_lower"].startswith(btn_text[:8])),
                     None)
        if match:
            btn["pptx_pos"] = {
                "left": match["left"],
                "top":  match["top"],
                "w":    match["w"],
                "h":    match["h"],
            }
            n_matched += 1
    return n_matched


def patch_html(html_path, pptx_path):
    pptx_buttons = extract_pptx_button_positions(pptx_path)
    with open(html_path, "r", encoding="utf-8") as f:
        html = f.read()
    m = SLIDES_RE.search(html)
    if not m:
        return f"FAIL  {os.path.basename(html_path)}: no SLIDES"
    slides = json.loads(m.group(1))

    total_matched = 0
    total_buttons = 0
    slides_fully_matched = 0
    for s in slides:
        total_buttons += len(s.get("buttons") or [])
        n = match_button(s.get("id"), slides, pptx_buttons)
        total_matched += n
        if n == len(s.get("buttons") or []) and n > 0:
            slides_fully_matched += 1

    new_json = json.dumps(slides, ensure_ascii=False, separators=(", ", ": "))
    new_html = html[:m.start(1)] + new_json + html[m.end(1):]

    # Inject engine block (only if not already present)
    if ENGINE_MARKER not in new_html:
        if "</body>" in new_html:
            new_html = new_html.replace("</body>", ENGINE_BLOCK + "\n</body>", 1)

    with open(html_path, "w", encoding="utf-8") as f:
        f.write(new_html)
    return f"OK    {os.path.basename(html_path)}: {total_matched}/{total_buttons} buttons positioned, {slides_fully_matched} slides fully matched"


def main():
    print("Extracting PPTX button positions and patching HTMLs\n")
    for fname, pptx_path in PPTX_MAP.items():
        path = os.path.join(ROOT, fname)
        if not os.path.exists(path):
            print(f"MISS  {fname}")
            continue
        if not os.path.exists(pptx_path):
            print(f"FAIL  {fname}: PPTX not found")
            continue
        print(patch_html(path, pptx_path))

    # Mirror to drake_classroom (uses drake PPTX)
    dc = os.path.join(ROOT, "drake_classroom.html")
    if os.path.exists(dc):
        print(patch_html(dc, PPTX_MAP["drake.html"]))


if __name__ == "__main__":
    main()
