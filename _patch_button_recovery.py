"""Recover all missing PPTX buttons by reading hyperlinks from the source PPTX.

Original converter logic dropped buttons when its position-matching between
text-shapes and invisible click-target rectangles wasn't tight enough. This
script does a more aggressive match: for each "button-like" text shape on a
slide, find the NEAREST click-target rect and use its hyperlink target.

Adds any newly-recovered buttons (with their pptx_pos for hotspot rendering)
to the SLIDES JSON. Skips buttons that already exist in the JSON.
"""

import os
import re
import json
import math
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.abspath(__file__))

PPTX_MAP = {
    "columbus.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
    "cortereal.html": r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Corte-Real) FINAL.pptx",
    "dagama.html":    r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Da Gama) FINAL.pptx",
    "drake.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Drake) FINAL.pptx",
    "magellan.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Magellan) FINAL.pptx",
    "narvaez.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Narvaez) FINAL.pptx",
    "raleigh.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Raleigh) FINAL.pptx",
    "index.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
}

SLIDES_RE = re.compile(r"const SLIDES = (\[.*?\]);", re.DOTALL)


def find_click_targets_for_slide(slide, all_slides_list):
    """Return list of click-target rects: {x1,y1,x2,y2,target_idx} as %."""
    sw = slide.part.package.presentation_part.presentation.slide_width
    sh = slide.part.package.presentation_part.presentation.slide_height
    slide_part = slide.part
    targets = []
    cur_idx = next((i for i, s in enumerate(all_slides_list) if s is slide), -1)

    for shape in slide.shapes:
        # Each shape can have hlinkClick(s); shapes are usually invisible click-target rects
        hlinks = shape._element.findall('.//' + qn('a:hlinkClick'))
        if not hlinks:
            continue
        for h in hlinks:
            rid = h.get(qn('r:id'))
            action = h.get('action', '') or ''
            target_idx = None
            if rid and rid in slide_part.rels:
                rel = slide_part.rels[rid]
                # Skip external (URL) hyperlinks — only follow slide jumps
                try:
                    if rel.is_external:
                        continue
                    tp = rel.target_part
                    for j, s in enumerate(all_slides_list):
                        if s.part is tp:
                            target_idx = j; break
                except (ValueError, AttributeError):
                    continue
            elif 'nextslide' in action:
                target_idx = cur_idx + 1 if cur_idx + 1 < len(all_slides_list) else None
            elif 'previousslide' in action:
                target_idx = cur_idx - 1 if cur_idx > 0 else None
            elif 'firstslide' in action:
                target_idx = 0
            elif 'lastslide' in action:
                target_idx = len(all_slides_list) - 1
            if target_idx is None or target_idx == cur_idx:
                continue  # skip self-loops
            try:
                left = shape.left / sw * 100
                top  = shape.top  / sh * 100
                w    = shape.width / sw * 100
                h_   = shape.height / sh * 100
                targets.append({
                    "x1": left, "y1": top,
                    "x2": left + w, "y2": top + h_,
                    "cx": left + w/2, "cy": top + h_/2,
                    "target_idx": target_idx,
                    "shape_name": shape.name,
                })
            except Exception:
                pass
    return targets


def find_text_buttons(slide, presentation):
    """Return list of button-like text shapes: {text, left,top,w,h cx,cy} as %."""
    sw = presentation.slide_width
    sh = presentation.slide_height
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        text = shape.text_frame.text.strip()
        if not text: continue
        try:
            w_pct = shape.width / sw * 100
            h_pct = shape.height / sh * 100
            # Heuristic for "button": width < 80%, height < 30%
            if w_pct > 80 or h_pct > 30: continue
            l = shape.left / sw * 100
            t = shape.top / sh * 100
            out.append({
                "text": text,
                "left": l, "top": t,
                "w": w_pct, "h": h_pct,
                "cx": l + w_pct/2, "cy": t + h_pct/2,
            })
        except Exception:
            pass
    return out


def match_text_to_click(text_btns, click_targets):
    """Match each text button to the nearest click-target rect (greedy by distance).
    A click-target's center should be inside the text shape, OR vice-versa, OR within
    a few % of distance.
    """
    matches = {}  # idx in text_btns -> click_target dict
    used_clicks = set()
    # Try strongest match first: click-target center INSIDE text shape (most likely click-overlay-on-text)
    for ti, t in enumerate(text_btns):
        for ci, c in enumerate(click_targets):
            if ci in used_clicks: continue
            if t["left"] <= c["cx"] <= t["left"] + t["w"] and t["top"] <= c["cy"] <= t["top"] + t["h"]:
                matches[ti] = c
                used_clicks.add(ci)
                break
    # Fill remaining: nearest by center distance
    for ti, t in enumerate(text_btns):
        if ti in matches: continue
        best_ci, best_d = None, 1e9
        for ci, c in enumerate(click_targets):
            if ci in used_clicks: continue
            d = math.hypot(t["cx"] - c["cx"], t["cy"] - c["cy"])
            if d < best_d:
                best_d, best_ci = d, ci
        if best_ci is not None and best_d < 30:  # within 30% slide-distance
            matches[ti] = click_targets[best_ci]
            used_clicks.add(best_ci)
    return matches


def patch_html(html_path, pptx_path):
    p = Presentation(pptx_path)
    all_slides = list(p.slides)

    with open(html_path, "r", encoding="utf-8") as f:
        html = f.read()
    m = SLIDES_RE.search(html)
    slides_json = json.loads(m.group(1))

    n_added = 0
    n_replaced = 0
    for slide_idx, slide in enumerate(all_slides):
        slide_id = f"slide_{slide_idx}"
        s = next((x for x in slides_json if x.get("id") == slide_id), None)
        if not s: continue

        click_targets = find_click_targets_for_slide(slide, all_slides)
        text_btns = find_text_buttons(slide, p)
        matches = match_text_to_click(text_btns, click_targets)

        # Build new buttons list from matches
        new_buttons = []
        for ti, c in matches.items():
            t = text_btns[ti]
            target_id = f"slide_{c['target_idx']}"
            new_buttons.append({
                "text": t["text"],
                "target": target_id,
                "pptx_pos": {
                    "left": round(t["left"], 1),
                    "top":  round(t["top"], 1),
                    "w":    round(t["w"], 1),
                    "h":    round(t["h"], 1),
                },
            })

        if not new_buttons:
            continue  # leave existing buttons untouched

        # Preserve existing button properties (e.g. our v2 patches)
        existing_by_text = { (b.get("text") or "").strip().lower()[:30]: b
                             for b in (s.get("buttons") or []) }
        merged = []
        for nb in new_buttons:
            key = nb["text"].strip().lower()[:30]
            if key in existing_by_text:
                ex = existing_by_text[key]
                # Update target/position from PPTX, keep any other props
                ex["target"] = nb["target"]
                ex["pptx_pos"] = nb["pptx_pos"]
                ex["text"] = nb["text"]  # use full text from PPTX
                merged.append(ex)
                n_replaced += 1
            else:
                merged.append(nb)
                n_added += 1

        # Preserve any synthesized buttons (Continue / Back to Hub) we added that
        # have no PPTX equivalent — keep them at the end as a fallback.
        for ex in (s.get("buttons") or []):
            txt = (ex.get("text") or "").strip()
            if txt in ("Continue", "Back to Hub"):
                if not any(b.get("text") == txt for b in merged):
                    merged.append(ex)

        s["buttons"] = merged

    new_json = json.dumps(slides_json, ensure_ascii=False, separators=(", ", ": "))
    new_html = html[:m.start(1)] + new_json + html[m.end(1):]
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(new_html)
    return f"OK    {os.path.basename(html_path)}: +{n_added} new, {n_replaced} updated"


def main():
    print("Recovering missing PPTX buttons across all character HTMLs\n")
    for fname, pptx_path in PPTX_MAP.items():
        html_path = os.path.join(ROOT, fname)
        if not os.path.exists(html_path):
            print(f"MISS  {fname}")
            continue
        if not os.path.exists(pptx_path):
            print(f"FAIL  {fname}: PPTX not found")
            continue
        print(patch_html(html_path, pptx_path))

    dc = os.path.join(ROOT, "drake_classroom.html")
    if os.path.exists(dc):
        print(patch_html(dc, PPTX_MAP["drake.html"]))


if __name__ == "__main__":
    main()
