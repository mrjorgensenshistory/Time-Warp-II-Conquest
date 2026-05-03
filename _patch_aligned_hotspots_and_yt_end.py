"""Two integrated, idempotent patches:

1. ALIGNED PPTX HOTSPOTS — adds `pptx_pos` (left/top/w/h %) to existing buttons
   in SLIDES JSON. DOES NOT rewrite targets (that was the bug from earlier today).
   Plus injects a CSS wrapper that pins the slide background to a 16:9 box
   regardless of viewport, so hotspot percentages align EXACTLY with the
   PPTX image. Hides engine bottom-row buttons only on slides where every
   button has a position (so partial-coverage slides still get a fallback).

2. YT END-OF-SONG bg-music RESTORE — replaces the simple iframe with one
   that uses the YouTube IFrame Player API. When the embedded video ends
   naturally, bg-music unmutes automatically. (Previously bg-music only
   unmuted when the player navigated away, leaving silence if the song
   ended first.)

All changes are idempotent — markers prevent double-application.
"""

import os
import re
import json
from pptx import Presentation

ROOT = os.path.dirname(os.path.abspath(__file__))

PPTX_MAP = {
    "columbus.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
    "cortereal.html": r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Corte-Real) FINAL.pptx",
    "dagama.html":    r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Da Gama) FINAL.pptx",
    "drake.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Drake) FINAL.pptx",
    "magellan.html":  r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Magellan) FINAL.pptx",
    "narvaez.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Narvaez) FINAL.pptx",
    "raleigh.html":   r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Raleigh) FINAL.pptx",
    "index.html":     r"F:\Michael's\ACE\Teach With Magic\Experiences\Time Warp\Explorers\Time Warp II_ Conquest (Tutorial w Columbus) FINAL.pptx",
}
ALL_HTMLS = list(PPTX_MAP.keys()) + ["drake_classroom.html"]
SLIDES_RE = re.compile(r"const SLIDES = (\[.*?\]);", re.DOTALL)

# === Phase 1: Position annotations (no target rewriting) ===

def extract_pptx_button_positions(pptx_path):
    """Returns dict: page_index_0based -> [{text, text_lower, left%, top%, w%, h%}]"""
    p = Presentation(pptx_path)
    sw, sh = p.slide_width, p.slide_height
    out = {}
    for i, slide in enumerate(p.slides):
        page_btns = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            try:
                w_pct = shape.width / sw * 100
                h_pct = shape.height / sh * 100
                # Heuristic for button: not full-slide (<80% wide, <30% tall)
                if w_pct > 80 or h_pct > 35:
                    continue
                page_btns.append({
                    "text": text,
                    "text_lower": text.lower(),
                    "left": round(shape.left / sw * 100, 1),
                    "top":  round(shape.top  / sh * 100, 1),
                    "w":    round(w_pct, 1),
                    "h":    round(h_pct, 1),
                })
            except Exception:
                pass
        out[i] = page_btns
    return out


def annotate_positions(html_path, pptx_path):
    """Add pptx_pos to existing buttons. NEVER changes target."""
    pptx_btns = extract_pptx_button_positions(pptx_path)
    with open(html_path, "r", encoding="utf-8") as f:
        html = f.read()
    m = SLIDES_RE.search(html)
    if not m:
        return f"FAIL  {os.path.basename(html_path)}: no SLIDES"
    slides = json.loads(m.group(1))

    n_pos = 0
    n_buttons = 0
    fully_matched = 0
    for s in slides:
        sid = s.get("id", "")
        try:
            page = int(sid.split("_")[1])
        except Exception:
            continue
        page_btns = pptx_btns.get(page, [])
        slide_btns = s.get("buttons") or []
        n_buttons += len(slide_btns)
        slide_n_pos = 0
        for btn in slide_btns:
            btn_text = (btn.get("text") or "").lower().strip()
            if not btn_text:
                continue
            # Match by prefix (PPTX text often has more chars than the truncated HTML text)
            prefix = btn_text[:15]
            short = btn_text[:8]
            match = next((p for p in page_btns
                          if (prefix and prefix in p["text_lower"]) or
                             (short and p["text_lower"].startswith(short))),
                         None)
            if match:
                btn["pptx_pos"] = {"left": match["left"], "top": match["top"],
                                   "w": match["w"], "h": match["h"]}
                n_pos += 1
                slide_n_pos += 1
        if slide_btns and slide_n_pos == len(slide_btns):
            fully_matched += 1

    new_json = json.dumps(slides, ensure_ascii=False, separators=(", ", ": "))
    new_html = html[:m.start(1)] + new_json + html[m.end(1):]
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(new_html)
    return f"OK    {os.path.basename(html_path)}: {n_pos}/{n_buttons} buttons annotated, {fully_matched} slides fully matched"


# === Phase 2: 16:9 aspect wrapper + hotspot rendering engine ===

ENGINE_MARKER = "<!-- TW2-ALIGNED-HOTSPOTS-V1 -->"

ENGINE_BLOCK = """
<!-- TW2-ALIGNED-HOTSPOTS-V1 -->
<style>
/* 16:9 aspect-ratio wrapper that holds the slide bg + hotspots together.
   Centered in the viewport, sized to whichever dimension limits the 16:9 fit.
   Hotspots positioned inside this wrap with PPTX % become exactly aligned. */
#tw2-aspect-frame {
  position: absolute;
  top: 50%; left: 50%;
  transform: translate(-50%, -50%);
  width: 100vw;
  height: calc(100vw * 9 / 16);
  max-height: 100vh;
  aspect-ratio: 16 / 9;
  pointer-events: none;
  z-index: 1;
}
@media (min-aspect-ratio: 16/9) {
  #tw2-aspect-frame {
    width: calc(100vh * 16 / 9);
    height: 100vh;
  }
}
.tw2-hotspot {
  position: absolute;
  z-index: 16;
  cursor: pointer;
  border: 2px dashed rgba(201, 161, 74, 0.55);
  background: rgba(201, 161, 74, 0.06);
  border-radius: 4px;
  pointer-events: auto;
  transition: all 0.15s ease;
  box-shadow: 0 0 12px rgba(201,161,74,0.15) inset;
}
.tw2-hotspot:hover {
  border-color: rgba(255, 215, 0, 1);
  border-style: solid;
  background: rgba(201, 161, 74, 0.2);
  box-shadow: 0 0 18px rgba(255, 215, 0, 0.45);
  transform: scale(1.02);
}
.tw2-hotspot:active { background: rgba(201, 161, 74, 0.4); }
@keyframes tw2-hotspot-pulse {
  0%, 100% { box-shadow: 0 0 12px rgba(201,161,74,0.12) inset, 0 0 6px rgba(201,161,74,0.3); }
  50%      { box-shadow: 0 0 12px rgba(201,161,74,0.25) inset, 0 0 14px rgba(255,215,0,0.55); }
}
.tw2-hotspot { animation: tw2-hotspot-pulse 2.4s ease-in-out infinite; }

/* Hide the engine bottom-button row on slides where every button is positioned */
.slide.tw2-fully-positioned .buttons {
  display: none !important;
}
</style>
<script>
(function() {
  if (typeof window.goToSlide !== 'function') return;
  // Ensure aspect-frame container exists once
  function ensureFrame() {
    var f = document.getElementById('tw2-aspect-frame');
    if (!f) {
      f = document.createElement('div');
      f.id = 'tw2-aspect-frame';
      var game = document.getElementById('game') || document.body;
      game.appendChild(f);
    }
    return f;
  }
  var prev = window.goToSlide;
  window.goToSlide = function(id) {
    prev(id);
    if (typeof id !== 'string' || id === 'HOME' || id.endsWith('.html')) return;
    if (typeof SLIDES === 'undefined') return;
    var slide = SLIDES.find(function(s){ return s.id === id; });
    if (!slide || !slide.buttons || !slide.buttons.length) return;

    setTimeout(function() {
      var frame = ensureFrame();
      // Clear previous hotspots
      frame.innerHTML = '';
      var positioned = 0;
      slide.buttons.forEach(function(btn) {
        if (!btn.pptx_pos) return;
        var hs = document.createElement('div');
        hs.className = 'tw2-hotspot';
        hs.style.left   = btn.pptx_pos.left + '%';
        hs.style.top    = btn.pptx_pos.top  + '%';
        hs.style.width  = btn.pptx_pos.w    + '%';
        hs.style.height = btn.pptx_pos.h    + '%';
        hs.title = btn.text;
        hs.onclick = function() { window.goToSlide(btn.target); };
        frame.appendChild(hs);
        positioned++;
      });
      var active = document.querySelector('.slide.active');
      if (active && positioned === slide.buttons.length && positioned > 0) {
        active.classList.add('tw2-fully-positioned');
      }
    }, 80);
  };
})();
</script>"""


def inject_engine(path):
    with open(path, "r", encoding="utf-8") as f:
        html = f.read()
    if ENGINE_MARKER in html:
        return f"SKIP  {os.path.basename(path)} (already patched)"
    if "</body>" not in html:
        return f"FAIL  {os.path.basename(path)}: no </body>"
    new_html = html.replace("</body>", ENGINE_BLOCK + "\n</body>", 1)
    with open(path, "w", encoding="utf-8") as f:
        f.write(new_html)
    return f"OK    {os.path.basename(path)}: aligned-hotspot engine injected"


# === Phase 3: YT player end-of-song bg-music restore ===

YT_END_MARKER = "<!-- TW2-YT-END-RESTORE-V1 -->"
YT_END_BLOCK = """
<!-- TW2-YT-END-RESTORE-V1 -->
<script>
(function() {
  // Load YouTube IFrame Player API once
  if (!document.getElementById('yt-iframe-api-script')) {
    var tag = document.createElement('script');
    tag.id = 'yt-iframe-api-script';
    tag.src = 'https://www.youtube.com/iframe_api';
    document.head.appendChild(tag);
  }
  // Track active player so we can listen for end-of-song
  window._tw2_yt_player = null;

  function bindEndListener() {
    var iframes = document.querySelectorAll('.yt-overlay iframe');
    iframes.forEach(function(iframe) {
      // Ensure enablejsapi=1 so the API can attach
      try {
        var src = iframe.src || '';
        if (src && src.indexOf('enablejsapi=1') === -1 && src.indexOf('about:blank') === -1) {
          iframe.src = src + (src.indexOf('?') >= 0 ? '&' : '?') + 'enablejsapi=1';
        }
      } catch(e) {}
      if (!window.YT || !window.YT.Player) return;
      if (iframe._tw2_bound) return;
      iframe._tw2_bound = true;
      try {
        window._tw2_yt_player = new window.YT.Player(iframe, {
          events: {
            'onStateChange': function(event) {
              // YT.PlayerState.ENDED === 0
              if (event.data === 0) {
                var bg = document.getElementById('bg-music');
                if (bg) {
                  bg.muted = false;
                  if (bg.paused) {
                    try { bg.volume = 0.45; bg.play().catch(function(){}); } catch(e){}
                  }
                }
              }
            }
          }
        });
      } catch(e) {}
    });
  }

  // Watch the game div for new yt-overlay nodes (engine creates them async)
  var observer = new MutationObserver(function(muts) {
    var hasYt = false;
    muts.forEach(function(m) {
      m.addedNodes && m.addedNodes.forEach(function(n) {
        if (n.nodeType === 1 && (n.classList && n.classList.contains('yt-overlay') ||
            (n.querySelector && n.querySelector('.yt-overlay')))) {
          hasYt = true;
        }
      });
    });
    if (hasYt) {
      // Slight delay so iframe src is set first by the existing engine
      setTimeout(bindEndListener, 400);
    }
  });
  document.addEventListener('DOMContentLoaded', function() {
    var game = document.getElementById('game');
    if (game) observer.observe(game, { childList: true, subtree: true });
    // Also handle YouTube API readiness — bind any iframes already present
    window.onYouTubeIframeAPIReady = function() { bindEndListener(); };
  });
})();
</script>"""


def inject_yt_end_restore(path):
    with open(path, "r", encoding="utf-8") as f:
        html = f.read()
    if YT_END_MARKER in html:
        return f"SKIP  {os.path.basename(path)} (yt-end already patched)"
    if "</body>" not in html:
        return f"FAIL  {os.path.basename(path)}: no </body>"
    new_html = html.replace("</body>", YT_END_BLOCK + "\n</body>", 1)
    with open(path, "w", encoding="utf-8") as f:
        f.write(new_html)
    return f"OK    {os.path.basename(path)}: yt-end-restore injected"


def main():
    print("=== Phase 1: Annotate buttons with PPTX positions ===\n")
    for fname, pptx_path in PPTX_MAP.items():
        path = os.path.join(ROOT, fname)
        if not os.path.exists(path) or not os.path.exists(pptx_path):
            print(f"MISS  {fname}")
            continue
        print(annotate_positions(path, pptx_path))
    # drake_classroom mirrors drake
    dc = os.path.join(ROOT, "drake_classroom.html")
    if os.path.exists(dc):
        print(annotate_positions(dc, PPTX_MAP["drake.html"]))

    print()
    print("=== Phase 2: Inject 16:9 aspect-frame + hotspot engine ===\n")
    for fname in ALL_HTMLS:
        path = os.path.join(ROOT, fname)
        if not os.path.exists(path): continue
        print(inject_engine(path))

    print()
    print("=== Phase 3: Inject YT end-of-song bg-music restore ===\n")
    for fname in ALL_HTMLS:
        path = os.path.join(ROOT, fname)
        if not os.path.exists(path): continue
        print(inject_yt_end_restore(path))


if __name__ == "__main__":
    main()
